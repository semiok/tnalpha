"""文档抽文本（复用 tngen extract 逻辑）。pdf/docx/pptx/xlsx → 纯文本。

上传时调 `extract_text` 把文档内容抽出来存库，供 AI 解析读真实内容。
容错：不支持的格式 / 解析失败 → 返回空串，不抛（不拖垮上传）。
"""
from pathlib import Path


def extract_text(path: str) -> str:
    """按扩展名抽文本。失败或不支持返回空串。"""
    ft = Path(path).suffix.lower().lstrip(".")
    try:
        if ft == "pdf":
            return _pdf(path)
        if ft == "docx":
            return _docx(path)
        if ft == "pptx":
            return _pptx(path)
        if ft == "xlsx":
            return _xlsx(path)
    except Exception:
        return ""  # 解析失败不拖垮上传，正文留空
    return ""      # 不支持的格式（如 txt/图片）留空


def _pdf(path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(path)
    return "\n".join((p.extract_text() or "") for p in reader.pages).strip()


def _docx(path: str) -> str:
    from docx import Document
    return "\n".join(p.text for p in Document(path).paragraphs).strip()


def _pptx(path: str) -> str:
    from pptx import Presentation
    out = []
    for slide in Presentation(path).slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                out.append(shape.text_frame.text)
    return "\n".join(out).strip()


def _xlsx(path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            out.append(" ".join("" if c is None else str(c) for c in row))
    return "\n".join(out).strip()
